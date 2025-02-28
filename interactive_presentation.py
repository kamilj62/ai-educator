import tkinter as tk
from tkinter import ttk, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import openai
from dotenv import load_dotenv
import os
import json

class PresentationCreator:
    def __init__(self):
        load_dotenv()
        self.openai_api_key = os.getenv('OPENAI_API_KEY')
        if not self.openai_api_key:
            raise ValueError("Please set your OpenAI API key in the .env file")
        
        self.client = openai.OpenAI(api_key=self.openai_api_key)
        
        self.root = tk.Tk()
        self.root.title("AI-Powered Presentation Creator")
        self.root.geometry("800x600")
        
        self.slides = []
        self.current_slide = 0
        
        self.setup_ui()
    
    def setup_ui(self):
        # Keywords input
        keyword_frame = ttk.Frame(self.root)
        keyword_frame.pack(pady=10, padx=10, fill="x")
        
        ttk.Label(keyword_frame, text="Enter Keywords (comma-separated):").pack(side="left")
        self.keyword_entry = ttk.Entry(keyword_frame, width=50)
        self.keyword_entry.pack(side="left", padx=5)
        
        # Number of slides
        slides_frame = ttk.Frame(self.root)
        slides_frame.pack(pady=5, padx=10, fill="x")
        
        ttk.Label(slides_frame, text="Number of Slides:").pack(side="left")
        self.slides_spinbox = ttk.Spinbox(slides_frame, from_=1, to=10, width=5)
        self.slides_spinbox.set(5)
        self.slides_spinbox.pack(side="left", padx=5)
        
        # Preview area
        preview_frame = ttk.LabelFrame(self.root, text="Preview")
        preview_frame.pack(pady=10, padx=10, fill="both", expand=True)
        
        self.preview_text = tk.Text(preview_frame, height=15, wrap="word")
        self.preview_text.pack(fill="both", expand=True, padx=5, pady=5)
        
        # Buttons
        button_frame = ttk.Frame(self.root)
        button_frame.pack(pady=10, padx=10)
        
        ttk.Button(button_frame, text="Generate Content", command=self.generate_content).pack(side="left", padx=5)
        ttk.Button(button_frame, text="Save Presentation", command=self.save_presentation).pack(side="left", padx=5)
        
        # Status label
        self.status_label = ttk.Label(self.root, text="")
        self.status_label.pack(pady=5)
    
    def generate_content(self):
        keywords = [k.strip() for k in self.keyword_entry.get().split(',') if k.strip()]
        if not keywords:
            messagebox.showerror("Error", "Please enter at least one keyword!")
            return
        
        num_slides = int(self.slides_spinbox.get())
        
        try:
            prompt = f"""Create a presentation outline with {num_slides} slides based on these keywords: {', '.join(keywords)}.
            For each slide, provide:
            1. A clear, concise title
            2. 3-4 key bullet points
            Format the response as JSON with this structure:
            {{
                "slides": [
                    {{"title": "Slide Title", "content": ["Point 1", "Point 2", "Point 3"]}}
                ]
            }}"""
            
            response = self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7
            )
            
            content = response.choices[0].message.content
            presentation_data = json.loads(content)
            
            self.slides = presentation_data["slides"]
            
            # Update preview
            self.preview_text.delete("1.0", "end")
            for i, slide in enumerate(self.slides, 1):
                self.preview_text.insert("end", f"\nSlide {i}: {slide['title']}\n")
                for point in slide['content']:
                    self.preview_text.insert("end", f"• {point}\n")
                self.preview_text.insert("end", "\n")
            
            self.status_label.config(text="Content generated successfully!")
            
        except Exception as e:
            messagebox.showerror("Error", f"Failed to generate content: {str(e)}")
    
    def save_presentation(self):
        if not self.slides:
            messagebox.showerror("Error", "Generate content first!")
            return
            
        prs = Presentation()
        
        for slide_data in self.slides:
            slide_layout = prs.slide_layouts[1]  # Using title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data["title"]
            
            # Set content
            content = slide.placeholders[1]
            content_text = content.text_frame
            
            for point in slide_data["content"]:
                p = content_text.add_paragraph()
                p.text = point
                p.level = 0
            
        prs.save("ai_generated_presentation.pptx")
        self.status_label.config(text="Presentation saved as 'ai_generated_presentation.pptx'!")
    
    def run(self):
        self.root.mainloop()

if __name__ == "__main__":
    app = PresentationCreator()
    app.run()
