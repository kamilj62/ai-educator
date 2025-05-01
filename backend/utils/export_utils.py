"""Utility functions for exporting presentations with secure image handling."""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

def create_presentation(title: str, slides: list, output_dir: str = "exports") -> str:
    """Create a PowerPoint presentation with secure image handling."""
    # Create exports directory if it doesn't exist
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    # Initialize presentation
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    # Add content slides
    for slide_content in slides:
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Add title
        title = content_slide.shapes.title
        title.text = slide_content["title"]
        
        # Add bullet points
        body_shape = content_slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for point in slide_content["bullet_points"]:
            p = tf.add_paragraph()
            p.text = point["text"]
            p.level = 0
        
        # Add image if present
        if "image_url" in slide_content and slide_content["image_url"]:
            img_path = os.path.join("marvelAI", slide_content["image_url"])
            if os.path.exists(img_path):
                content_slide.shapes.add_picture(
                    img_path,
                    Inches(1),
                    Inches(3),
                    width=Inches(8)
                )
        
        # Add examples
        if slide_content.get("examples"):
            p = tf.add_paragraph()
            p.text = "\nExamples:"
            p.font.bold = True
            for example in slide_content["examples"]:
                p = tf.add_paragraph()
                p.text = "• " + example["text"]
                p.level = 1
        
        # Add discussion questions
        if slide_content.get("discussion_questions"):
            p = tf.add_paragraph()
            p.text = "\nDiscussion Questions:"
            p.font.bold = True
            for question in slide_content["discussion_questions"]:
                p = tf.add_paragraph()
                p.text = "• " + question
                p.level = 1
    
    # Save with timestamp
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{title.replace(' ', '_')}_{timestamp}.pptx"
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)
    
    return output_path
