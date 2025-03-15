from typing import List, Optional
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from models import SlideContent, ExportFormat, Presentation, ExportRequest
import os
import requests
from io import BytesIO
from urllib.parse import urlparse
import logging

logger = logging.getLogger(__name__)

class PresentationService:
    def __init__(self, ai_service):
        """Initialize PresentationService with required dependencies."""
        self.ai_service = ai_service
        logger.info("PresentationService initialized with AIService")

    def _add_image_to_slide(self, slide, image_url: str, image_caption: Optional[str] = None):
        """Add an image to a slide with optional caption."""
        try:
            # Determine if URL or local path
            is_url = bool(urlparse(image_url).scheme)
            
            if is_url:
                # Download image from URL
                response = requests.get(image_url)
                response.raise_for_status()
                image_stream = BytesIO(response.content)
            else:
                # Open local file
                with open(image_url, 'rb') as img_file:
                    image_stream = BytesIO(img_file.read())
            
            # Add image to slide
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)  # Adjust size as needed
            pic = slide.shapes.add_picture(image_stream, left, top, width=width)
            
            # Add caption if provided
            if image_caption:
                left = Inches(1)
                top = top + pic.height + Inches(0.2)
                width = Inches(8)
                height = Inches(0.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = image_caption
                
        except Exception as e:
            logger.error(f"Failed to add image to slide: {str(e)}")
            raise

    def create_pptx(self, presentation: Presentation, output_path: str) -> str:
        """Create a PowerPoint presentation from slide content."""
        try:
            prs = PPTXPresentation()
            
            # Create title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Space Exploration: Recent Achievements and Future Missions"
            subtitle.text = f"A {presentation.instructional_level.value.replace('_', ' ').title()} Level Presentation\nGenerated with MarvelAI"
            
            # Create content slides
            for slide_content in presentation.slides:
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                # Set title
                title = slide.shapes.title
                title.text = slide_content.title
                
                # Add bullet points
                content = slide.placeholders[1]
                tf = content.text_frame
                tf.clear()  # Clear any default text
                
                # Add bullet points
                for point in slide_content.bullet_points:
                    p = tf.add_paragraph()
                    p.text = point.text
                    p.level = 0
                
                # Add examples section if present
                if slide_content.examples:
                    p = tf.add_paragraph()
                    p.text = "Examples:"
                    p.level = 0
                    for example in slide_content.examples:
                        p = tf.add_paragraph()
                        p.text = example.text
                        p.level = 1
                
                # Add discussion questions section
                if slide_content.discussion_questions:
                    p = tf.add_paragraph()
                    p.text = "Discussion Questions:"
                    p.level = 0
                    for question in slide_content.discussion_questions:
                        p = tf.add_paragraph()
                        p.text = question
                        p.level = 1
                
                # Add image if provided
                if slide_content.image_url:
                    self._add_image_to_slide(slide, slide_content.image_url, slide_content.image_caption)
            
            # Save the presentation
            prs.save(output_path)
            return output_path
            
        except Exception as e:
            logger.error(f"Error creating PowerPoint presentation: {str(e)}")
            raise ValueError(f"Failed to create presentation: {str(e)}")

    def export_presentation(self, request: ExportRequest) -> str:
        """Export the presentation in the specified format."""
        try:
            base_path = "presentations"
            os.makedirs(base_path, exist_ok=True)
            
            output_path = os.path.join(base_path, f"presentation.{request.format.value}")
            
            if request.format == ExportFormat.PPTX:
                return self.create_pptx(request.presentation, output_path)
            elif request.format == ExportFormat.PDF:
                return self.create_pptx(request.presentation, output_path.replace('.pdf', '.pptx'))
            else:
                raise ValueError(f"Unsupported export format: {request.format}")
                
        except Exception as e:
            logger.error(f"Error exporting presentation: {str(e)}")
            raise ValueError(f"Failed to export presentation: {str(e)}")
