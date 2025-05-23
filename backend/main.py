from fastapi import FastAPI, HTTPException, Body, Request
from fastapi.responses import JSONResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
import uvicorn
import os
from dotenv import load_dotenv
import logging
import traceback
import json
import time
from datetime import datetime
from pptx import Presentation as pptx_Presentation
from pptx.util import Inches
from contextlib import asynccontextmanager
import asyncio
from models import (
    PresentationInput, 
    OutlineResponse, 
    Presentation, 
    ExportRequest,
    SlideTopic,
    InstructionalLevel,
    SlideContent,
    SlideGenerationRequest,
    ImageGenerationRequest,
    SlideLayout,
    SlideContentNew,
    SlideNew as Slide
)
from ai_service import AIService
from rate_limiter import RateLimiter
from exceptions import (
    ImageGenerationError,
    ContentSafetyError,
    ErrorType,
    ImageServiceProvider
)

# Configure logging with more detail
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Load OpenAI API key from credentials or environment
try:
    api_key = os.environ.get('OPENAI_API_KEY')
    if not api_key:
        with open('credentials/marvelai-imagen-sa-key.json', 'r') as f:
            credentials = json.load(f)
            api_key = credentials.get('openai_api_key')
    if not api_key:
        raise ValueError("OpenAI API key not found in environment or credentials file")
    logger.info(f"✅ OpenAI API Key Loaded: {api_key[:5]}...{api_key[-5:]}")
except Exception as e:
    logger.error(f"Error loading credentials: {str(e)}")
    raise

# Initialize services
ai_service = None
rate_limiter = RateLimiter()

# Ensure static directories exist
import pathlib
STATIC_DIR = pathlib.Path(__file__).parent / "static"
IMAGES_DIR = STATIC_DIR / "images"
STATIC_DIR.mkdir(exist_ok=True)
IMAGES_DIR.mkdir(exist_ok=True)
logger.info(f"Static directories created: {STATIC_DIR}, {IMAGES_DIR}")

# Replace deprecated startup event with lifespan context
from contextlib import asynccontextmanager

@asynccontextmanager
async def lifespan(app):
    global ai_service
    try:
        ai_service = AIService(rate_limiter)
        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
        await ai_service.initialize()
    except Exception as e:
        logger.error(f"Failed to initialize AI service for testing: {e}")
        logger.error(traceback.format_exc())
    yield

app = FastAPI(lifespan=lifespan)

# CORS setup
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://ai-educator-eight.vercel.app",
        "https://ai-educator-one.vercel.app",
        "https://ai-educator-nine.vercel.app",
        "https://ai-educator-iewndpp7z-kamilj62s-projects.vercel.app",
        "http://localhost:3000"
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

app.mount("/static", StaticFiles(directory="static"), name="static")

# Basic error handling middleware
@app.middleware("http")
async def log_requests(request, call_next):
    try:
        start_time = time.time()
        response = await call_next(request)
        process_time = time.time() - start_time
        logger.info(f"Request to {request.url.path} took {process_time:.2f} seconds")
        return response
    except Exception as e:
        logger.error(f"Error processing request: {str(e)}")
        raise

# Health check endpoint
@app.get("/api/health")
async def health_check():
    return {"status": "healthy"}

@app.post("/api/generate/outline")
async def generate_outline(input_data: PresentationInput) -> OutlineResponse:
    try:
        logger.info(f"Received outline generation request: {input_data}")
        if not input_data.context.strip():
            logging.warning("[DEBUG] Context is empty!")
            raise HTTPException(status_code=400, detail="Context cannot be empty")
        if not 1 <= input_data.num_slides <= 20:
            logging.warning(f"[DEBUG] Invalid num_slides: {input_data.num_slides}")
            raise HTTPException(status_code=400, detail="Number of slides must be between 1 and 20")
        try:
            response = await ai_service.generate_outline(
                input_data.context,
                input_data.num_slides,
                input_data.instructional_level,
                sensitive=False
            )
            return response
        except Exception as e:
            logging.error(f"[DEBUG] Error in generate_outline: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(
                status_code=500,
                detail=f"Error generating outline: {str(e)}"
            )
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"[DEBUG] Unexpected error in generate_outline endpoint: {str(e)}")
        logging.error(traceback.format_exc())
        raise HTTPException(
            status_code=500,
            detail=f"Unexpected error: {str(e)}"
        )

@app.post("/api/generate/slide")
async def generate_slide_content(request: SlideGenerationRequest):
    try:
        logger.info(f"Received slide generation request for topic: {request.topic.title}")
        logger.debug(f"Request data: {request.model_dump()}")
        try:
            slide_content = await ai_service.generate_slide_content(
                request.topic,
                request.instructional_level,
                request.layout
            )
            logger.info(f"Successfully generated content for topic: {request.topic.title}")
            logger.debug(f"Slide content: {slide_content.model_dump()}")
            return {
                "data": {
                    "slide": slide_content.model_dump(),
                    "warnings": []  # Add any warnings if needed
                }
            }
        except ImageGenerationError as e:
            logger.error(f"Image generation error: {e.message} ({e.error_type.value})")
            return JSONResponse(
                status_code=500,
                content={
                    "status": "error",
                    "message": e.message,
                    "error_type": e.error_type.value,
                    "service": e.service.value if e.service else None,
                    "retry_after": e.retry_after
                }
            )
        except ValueError as ve:
            logger.error(f"Value error in generate_slide_content: {str(ve)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=400, detail=str(ve))
        except Exception as e:
            logger.error(f"Error in generate_slide_content: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(
                status_code=500,
                detail=f"Error generating slide content: {str(e)}"
            )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error in generate_slide_content endpoint: {str(e)}")
        logger.error(traceback.format_exc())
        raise HTTPException(
            status_code=500,
            detail=f"Unexpected error: {str(e)}"
        )

@app.post("/api/generate-image")
async def generate_image(request: dict):
    try:
        logger.info(f"Received image generation request: {request}")
        try:
            # Ensure prompt is present and valid
            prompt = request.get('prompt') if isinstance(request, dict) else None
            if not prompt or not isinstance(prompt, str) or not prompt.strip():
                logger.error("Image generation request missing valid 'prompt'")
                return JSONResponse(
                    status_code=400,
                    content={
                        "status": "error",
                        "message": "Image generation request must include a non-empty 'prompt' field.",
                        "error_type": "INVALID_REQUEST"
                    }
                )
            image_url = await ai_service.generate_image(request)
            logger.info(f"Generated image URL: {image_url}")
            return {"image_url": image_url}
        except ImageGenerationError as e:
            logger.error(f"Image generation error: {e.message} ({e.error_type.value})")
            return JSONResponse(
                status_code=500,
                content={
                    "status": "error",
                    "message": e.message,
                    "error_type": e.error_type.value,
                    "service": e.service.value if e.service else None,
                    "retry_after": e.retry_after
                }
            )
        except Exception as e:
            logger.error(f"Error in generate_image: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(
                status_code=500,
                detail=f"Error generating image: {str(e)}"
            )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error in generate_image endpoint: {str(e)}")
        logger.error(traceback.format_exc())
        raise HTTPException(
            status_code=500,
            detail=f"Unexpected error: {str(e)}"
        )

@app.post("/api/generate/slides")
async def generate_slides(request: SlideGenerationRequest):
    try:
        logger.debug(f"Received slide request: {request.model_dump()}")
        topic = request.topic
        layout = request.layout or "title-bullets"
        response = {
            "title": topic.title,
            "subtitle": "",
            "body": topic.description or "",
            "key_points": topic.key_points,
            "image": {
                "url": getattr(topic, "image_url", None) or "",
                "alt": topic.image_prompt,
                "caption": getattr(topic, "image_caption", None) or "",
                "service": getattr(topic, "image_service", None) or "generated"
            }
        }
        return response
    except Exception as e:
        logger.error(f"Error generating slides: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail={
                "type": "API_ERROR",
                "message": f"Failed to generate slides: {str(e)}",
                "context": {"error": str(e)}
            }
        )

@app.post("/export")
async def export_presentation(request: ExportRequest):
    try:
        logger.info(f"Received export request for format: {request.format}")
        result = ai_service.export_presentation(request)
        return {"file_path": result}
    except Exception as e:
        logger.error(f"Error exporting presentation: {str(e)}")
        logger.error(traceback.format_exc())
        raise HTTPException(
            status_code=500,
            detail=f"Error exporting presentation: {str(e)}"
        )

@app.post("/export/quick")
async def quick_export(presentation: Presentation):
    try:
        export_dir = os.path.join("static", "exports")
        os.makedirs(export_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"renewable_energy_{timestamp}.pptx"
        output_path = os.path.join(export_dir, filename)
        prs = pptx_Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Renewable Energy"
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        for slide in presentation.slides:
            content_slide = prs.slides.add_slide(prs.slide_layouts[1])
            content_slide.shapes.title.text = slide.title
            body_shape = content_slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            for point in slide.key_points:
                p = tf.add_paragraph()
                p.text = point
            if slide.image_url:
                img_path = os.path.join("marvelAI", slide.image_url)
                if os.path.exists(img_path):
                    content_slide.shapes.add_picture(
                        img_path,
                        Inches(1),
                        Inches(2),
                        width=Inches(8)
                    )
        prs.save(output_path)
        return {"file_path": f"exports/{filename}"}
    except Exception as e:
        logger.error(f"Export error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

def is_sensitive_topic(context: str) -> bool:
    context_lower = context.lower()
    SENSITIVE_TOPICS = [
        "israeli palestinian conflict",
        "israel palestine",
        "gaza",
        "west bank",
        "middle east conflict"
    ]
    return any(topic in context_lower for topic in SENSITIVE_TOPICS)

if __name__ == "__main__":
    import os
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run("backend.main:app", host="0.0.0.0", port=port, reload=False)
